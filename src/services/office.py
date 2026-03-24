"""
Office Document Generation Services

Provides capabilities to CREATE and EDIT Office documents:
- DOCX: Create Word documents
- XLSX: Create Excel spreadsheets
- PPTX: Create PowerPoint presentations
"""

import os
from typing import Any

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def create_docx(path: str, content: str, title: str | None = None) -> str:
    """
    Create a Microsoft Word document.

    Args:
        path: Output path
        content: Main text content (paragraphs separated by newlines)
        title: Optional document title

    Returns:
        Success message or error
    """
    if not docx:
        return "Error: python-docx not installed."

    try:
        doc = docx.Document()

        if title:
            doc.add_heading(title, 0)

        for paragraph in content.split("\n"):
            if paragraph.strip():
                doc.add_paragraph(paragraph)

        doc.save(path)
        return f"Successfully created Word document at {path}"
    except Exception as e:
        return f"Error creating DOCX: {e}"


def create_xlsx(path: str, data: list[list[Any]], sheet_name: str = "Sheet1") -> str:
    """
    Create an Excel spreadsheet.

    Args:
        path: Output path
        data: List of rows, where each row is a list of values
        sheet_name: Name of the first sheet

    Returns:
        Success message or error
    """
    if not openpyxl:
        return "Error: openpyxl not installed."

    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = sheet_name

        for row in data:
            ws.append(row)

        wb.save(path)
        return f"Successfully created Excel file at {path}"
    except Exception as e:
        return f"Error creating XLSX: {e}"


def add_sheet_xlsx(path: str, data: list[list[Any]], sheet_name: str) -> str:
    """
    Add a new sheet to an existing Excel file.

    Args:
        path: Path to existing XLSX
        data: List of rows
        sheet_name: Name of the new sheet

    Returns:
        Success message or error
    """
    if not openpyxl:
        return "Error: openpyxl not installed."

    if not os.path.exists(path):
        return f"Error: File not found: {path}"

    try:
        wb = openpyxl.load_workbook(path)
        if sheet_name in wb.sheetnames:
            return f"Error: Sheet '{sheet_name}' already exists."

        ws = wb.create_sheet(sheet_name)
        for row in data:
            ws.append(row)

        wb.save(path)
        return f"Successfully added sheet '{sheet_name}' to {path}"
    except Exception as e:
        return f"Error modifying XLSX: {e}"


def create_pptx(path: str, slides_data: list[dict[str, str]]) -> str:
    """
    Create a PowerPoint presentation.

    Args:
        path: Output path
        slides_data: List of dicts, each containing 'title' and 'content' (bullet points)
                     Example: [{"title": "Slide 1", "content": "Point 1\nPoint 2"}]

    Returns:
        Success message or error
    """
    if not Presentation:
        return "Error: python-pptx not installed."

    try:
        prs = Presentation()

        for slide_info in slides_data:
            # Add a slide with title and content layout (layout index 1 usually)
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            title = slide.shapes.title
            title.text = slide_info.get("title", "Untitled Slide")

            # Set content
            content_text = slide_info.get("content", "")
            if content_text:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame
                tf.text = content_text

        prs.save(path)
        return f"Successfully created PowerPoint at {path}"
    except Exception as e:
        return f"Error creating PPTX: {e}"
