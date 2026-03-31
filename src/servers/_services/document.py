"""
Document Parsing Services

Provides format-aware document parsing for various file types:
- PDF: Text extraction with OCR fallback detection
- DOCX: Microsoft Word document parsing
- PPTX: PowerPoint presentation text extraction
- XLSX: Excel spreadsheet parsing
"""

try:
    import docx
except ImportError:
    docx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


def parse_pdf(path: str) -> str:
    """
    Extract text from PDF file.

    Args:
        path: Path to PDF file

    Returns:
        Extracted text or error message
    """
    if pdfplumber is None:
        return "[PDF Error: pdfplumber not installed]"

    text = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if not page_text or len(page_text.strip()) < 20:
                    page_text = f"[Page {i + 1} is likely scanned. OCR required.]"
                text.append(page_text)
        return "\n".join(text)
    except Exception as e:
        return f"[PDF Error: {e}]"


def parse_docx(path: str) -> str:
    """
    Extract text from Microsoft Word document.

    Args:
        path: Path to DOCX file

    Returns:
        Extracted text or error message
    """
    if docx is None:
        return "[DOCX Error: python-docx not installed]"

    try:
        doc = docx.Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except Exception as e:
        return f"[DOCX Error: {e}]"


def parse_pptx(path: str) -> str:
    """
    Extract text from PowerPoint presentation.

    Args:
        path: Path to PPTX file

    Returns:
        Extracted text or error message
    """
    if Presentation is None:
        return "[PPTX Error: python-pptx not installed]"

    try:
        prs = Presentation(path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        return f"[PPTX Error: {e}]"


def parse_xlsx(path: str) -> str:
    """
    Extract text from Excel spreadsheet.

    Args:
        path: Path to XLSX file

    Returns:
        Extracted text with sheet names or error message
    """
    if openpyxl is None:
        return "[XLSX Error: openpyxl not installed]"

    try:
        wb = openpyxl.load_workbook(path, data_only=True)
        text = []
        for sheet in wb.sheetnames:
            text.append(f"--- Sheet: {sheet} ---")
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text.append("\t".join([str(c) for c in row if c is not None]))
        return "\n".join(text)
    except Exception as e:
        return f"[XLSX Error: {e}]"
