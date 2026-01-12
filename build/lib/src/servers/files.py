import os
import glob
from pathlib import Path
from typing import List, Optional
from mcp.server.fastmcp import FastMCP

# Document Parsers
import pdfplumber
import docx
from pptx import Presentation
import openpyxl

mcp = FastMCP("PolymathFilesystem")

# 安全限制：允许访问的根目录
# 在真实部署中，这应该更严格
ALLOWED_ROOTS = [
    os.path.abspath("materials"),
    os.path.abspath("drafts"),
    os.getcwd()
]

def validate_path(path: str) -> str:
    """Validate that path is within the current working directory."""
    abs_path = os.path.abspath(path)
    base_dir = os.getcwd()
    
    # 如果 abs_path 不是以 base_dir 开头，说明试图越权访问
    if not abs_path.startswith(base_dir):
        raise PermissionError(f"Access Denied: Path {path} is outside of the workspace sandbox.")
    
    return abs_path

# --------------------------
# Parsers
# --------------------------
def parse_pdf(path: str) -> str:
    text = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text.append(page.extract_text() or "")
    return "\n".join(text)

def parse_docx(path: str) -> str:
    doc = docx.Document(path)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def parse_xlsx(path: str) -> str:
    wb = openpyxl.load_workbook(path, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        text.append(f"--- Sheet: {sheet} ---")
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            text.append("\t".join([str(c) for c in row if c is not None]))
    return "\n".join(text)

# --------------------------
# Tools
# --------------------------
@mcp.tool()
def list_directory(path: str) -> str:
    """List files and directories at the given path."""
    valid_path = validate_path(path)
    if not os.path.exists(valid_path):
        return "Error: Path not found."
    
    try:
        items = os.listdir(valid_path)
        return "\n".join(items)
    except Exception as e:
        return f"Error listing directory: {str(e)}"

@mcp.tool()
def read_file(path: str) -> str:
    """
    Intelligently read a file based on its extension.
    Supports: .txt, .md, .py, .json, .pdf, .docx, .pptx, .xlsx
    """
    valid_path = validate_path(path)
    if not os.path.exists(valid_path):
        return f"Error: File {path} not found."
    
    ext = os.path.splitext(path)[1].lower()
    
    try:
        if ext == ".pdf":
            return parse_pdf(valid_path)
        elif ext == ".docx":
            return parse_docx(valid_path)
        elif ext == ".pptx":
            return parse_pptx(valid_path)
        elif ext in [".xlsx", ".xls"]:
            return parse_xlsx(valid_path)
        else:
            # Default to text
            with open(valid_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
    except Exception as e:
        return f"Error reading file: {str(e)}"

@mcp.tool()
def write_file(path: str, content: str) -> str:
    """Write text content to a file. Overwrites if exists."""
    valid_path = validate_path(path)
    try:
        # Ensure parent dir exists
        os.makedirs(os.path.dirname(valid_path), exist_ok=True)
        with open(valid_path, "w", encoding="utf-8") as f:
            f.write(content)
        return f"Successfully wrote to {path}"
    except Exception as e:
        return f"Error writing file: {str(e)}"

if __name__ == "__main__":
    mcp.run()
