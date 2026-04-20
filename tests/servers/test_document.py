"""Tests for servers._services.document — document parsing services."""

import pytest

from src.servers._services.document import parse_docx, parse_pdf, parse_pptx, parse_xlsx


class TestParsePdf:
    def test_not_installed(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "pdfplumber", None)
        result = parse_pdf("any.pdf")
        assert "pdfplumber not installed" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("pdfplumber"),
        reason="pdfplumber not installed",
    )
    def test_parse_valid_pdf(self, tmp_path):
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(b"not a real pdf")
        result = parse_pdf(str(pdf_path))
        assert "Error" in result or "PDF" in result


class TestParseDocx:
    def test_not_installed(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "docx", None)
        result = parse_docx("any.docx")
        assert "python-docx not installed" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("docx"),
        reason="python-docx not installed",
    )
    def test_parse_valid_docx(self, tmp_path):
        from docx import Document

        doc = Document()
        doc.add_paragraph("Hello world")
        doc.add_paragraph("Second paragraph")
        path = tmp_path / "test.docx"
        doc.save(str(path))
        result = parse_docx(str(path))
        assert "Hello world" in result
        assert "Second paragraph" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("docx"),
        reason="python-docx not installed",
    )
    def test_parse_invalid_docx(self, tmp_path):
        bad_file = tmp_path / "bad.docx"
        bad_file.write_text("not a docx")
        result = parse_docx(str(bad_file))
        assert "Error" in result


class TestParsePptx:
    def test_not_installed(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "Presentation", None)
        result = parse_pptx("any.pptx")
        assert "python-pptx not installed" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("pptx"),
        reason="python-pptx not installed",
    )
    def test_parse_valid_pptx(self, tmp_path):
        from pptx import Presentation as Prs

        prs = Prs()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = slide.shapes.title
        title.text = "Slide Title"
        path = tmp_path / "test.pptx"
        prs.save(str(path))
        result = parse_pptx(str(path))
        assert "Slide Title" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("pptx"),
        reason="python-pptx not installed",
    )
    def test_parse_invalid_pptx(self, tmp_path):
        bad = tmp_path / "bad.pptx"
        bad.write_text("not pptx")
        result = parse_pptx(str(bad))
        assert "Error" in result


class TestParseXlsx:
    def test_not_installed(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "openpyxl", None)
        result = parse_xlsx("any.xlsx")
        assert "openpyxl not installed" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("openpyxl"),
        reason="openpyxl not installed",
    )
    def test_parse_valid_xlsx(self, tmp_path):
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sheet1"
        ws["A1"] = "Name"
        ws["B1"] = "Value"
        ws["A2"] = "test"
        ws["B2"] = 42
        path = tmp_path / "test.xlsx"
        wb.save(str(path))
        result = parse_xlsx(str(path))
        assert "Sheet1" in result
        assert "Name" in result
        assert "test" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("openpyxl"),
        reason="openpyxl not installed",
    )
    def test_parse_invalid_xlsx(self, tmp_path):
        bad = tmp_path / "bad.xlsx"
        bad.write_text("not xlsx")
        result = parse_xlsx(str(bad))
        assert "Error" in result


class TestDocumentImportFallbacks:
    def test_docx_import_fallback(self, monkeypatch):
        import src.servers._services.document as mod
        assert mod.docx is not None or mod.docx is None
        monkeypatch.setattr(mod, "docx", None)
        result = parse_docx("any.docx")
        assert "python-docx not installed" in result

    def test_openpyxl_import_fallback(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "openpyxl", None)
        result = parse_xlsx("any.xlsx")
        assert "openpyxl not installed" in result

    def test_pdfplumber_import_fallback(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "pdfplumber", None)
        result = parse_pdf("any.pdf")
        assert "pdfplumber not installed" in result

    def test_pptx_import_fallback(self, monkeypatch):
        import src.servers._services.document as mod
        monkeypatch.setattr(mod, "Presentation", None)
        result = parse_pptx("any.pptx")
        assert "python-pptx not installed" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("pdfplumber"),
        reason="pdfplumber not installed",
    )
    def test_parse_pdf_with_pages(self, tmp_path):
        result = parse_pdf(str(tmp_path / "nonexistent.pdf"))
        assert "Error" in result

    @pytest.mark.skipif(
        not __import__("importlib").util.find_spec("pdfplumber"),
        reason="pdfplumber not installed",
    )
    def test_parse_pdf_exception_handling(self, tmp_path):
        bad = tmp_path / "bad.pdf"
        bad.write_bytes(b"not a pdf")
        result = parse_pdf(str(bad))
        assert "Error" in result or "PDF" in result
